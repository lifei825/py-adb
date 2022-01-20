from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

# 按顺序排列生成月的周报
zb = [
    '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/慧企周报-12月第1周.pptx',
    '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/慧企周报-12月第2周.pptx',
    '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/慧企周报-12月第3周.pptx',
    '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/慧企周报-12月第4周.pptx',
    '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/慧企周报-12月第5周.pptx',
]

latest_zb = zb[-1]

# 用上个月月报作为模板
m = '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/test.pptx'
# 生成的新月报
new_m = '/Users/lifei/Library/Containers/com.tencent.xinWeChat/Data/Library/Application Support/\
com.tencent.xinWeChat/2.0b4.0.9/f77af8be9b77bf59843ef8bfa4b6f57d/Message/MessageTemp/6da2f04acf67b43f2a8c7e1b13730f89/\
File/12月/12月报-慧企.pptx'

month = 12
weeks = ['月第一周', '月第二周', '月第三周', '月第四周', '月第五周'][:len(zb)]

bg_list = []
zy_list = []
fb_list = []
ph_list = []
gz_list = []

# wwwfw=[[服务名,实例数]], wwwfw_pj=[project]
wwwfw = []
wwwfw_pj = []
wwht = []
wwht_pj = []
nwwfw = []
nwwfw_pj = []
nwht = []
nwht_pj = []

# wwwfw_cpu = [{"cpu_ysy": 0, "cpu_wsy": 0}], wwwfw_mem = [{"mem_ysy": 0, "mem_wsy": 0}]
wwwfw_cpu = []
wwwfw_mem = []
wwht_cpu = []
wwht_mem = []
nwwfw_cpu = []
nwwfw_mem = []
nwht_cpu = []
nwht_mem = []

# es charts
es_month = []
es_index = []
es_des = ''

latest_zb_prs = Presentation(latest_zb)


def main():
    for f in zb:
        global prs, m_prs
        prs = Presentation(f)
        m_prs = Presentation(m)
        # ppt页数
        # pri                                       nt(len(prs.slides))
        # 变更信息收集
        read_page(11)
        # 支撑发版收集
        read_page(12)
        # 资源权限管理收集
        read_page(13)
        # 配合及排障收集
        read_page(14)
        # 问题和告警
        read_page(15)
        # 集群资源使用情况 cpu mem
        read_page(19)
        read_page(20)
        read_page(21)
        read_page(22)

    # 获取最后一周周报信息
    read_latest_zb(19)
    read_latest_zb(20)
    read_latest_zb(21)
    read_latest_zb(22)
    read_latest_zb(23)
    read_latest_zb(24)

    # 运维工作统计
    write_m_page(7)
    # 变更写入月报
    write_m_page(8)
    # 资源管理写入月报
    write_m_page(9)
    # 配合及排障写入月报
    write_m_page(10)
    # 发版写入月报
    write_m_page(11)
    # 问题告警写入月报
    write_m_page(12)
    # 本月工作成果
    write_m_page(14)
    # 集群资源可分配容量
    write_m_page(18)
    # 外网微服务区资源使用情况
    write_m_page(19)
    write_m_page(20)
    write_m_page(21)
    write_m_page(22)
    # es chart
    write_m_page(23)


def read_page(n):
    bg = prs.slides[n-1]
    for i in bg.shapes:
        if i.has_chart:
            # 外网微服务cpu mem
            if n in [19, 20, 21, 22]:
                c = i.chart
                title = c.chart_title.text_frame.text
                if title.upper().find('CPU') >= 0:
                    cpu = [i.values for i in c.series][0]
                    if n == 19:
                        wwwfw_cpu.append({"cpu_ysy": cpu[0], "cpu_wsy": cpu[1]})
                    if n == 20:
                        wwht_cpu.append({"cpu_ysy": cpu[0], "cpu_wsy": cpu[1]})
                    if n == 21:
                        nwwfw_cpu.append({"cpu_ysy": cpu[0], "cpu_wsy": cpu[1]})
                    if n == 22:
                        nwht_cpu.append({"cpu_ysy": cpu[0], "cpu_wsy": cpu[1]})
                elif title.find("内存") >= 0:
                    mem = [i.values for i in c.series][0]
                    if n == 19:
                        wwwfw_mem.append({"mem_ysy": mem[0], "mem_wsy": mem[1]})
                    if n == 20:
                        wwht_mem.append({"mem_ysy": mem[0], "mem_wsy": mem[1]})
                    if n == 21:
                        nwwfw_mem.append({"mem_ysy": mem[0], "mem_wsy": mem[1]})
                    if n == 22:
                        nwht_mem.append({"mem_ysy": mem[0], "mem_wsy": mem[1]})

        if i.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = i.table
            cnt_row, cnt_col = len(tbl.rows), len(tbl.columns)
            for r in range(0, cnt_row):
                row_list = []
                for c in range(0, cnt_col):
                    cell = tbl.cell(r, c)
                    if len(cell.text) != 0:
                        # 变更
                        if n == 11:
                            if r > 0 and c == 2:
                                row_list.append(cell.text)

                        # 发版 or 问题告警
                        if n == 12 or n == 15:
                            if r > 0:
                                row_list.append(cell.text)
                        # 资源权限管理 or 配合操作及排障
                        if n == 13 or n == 14:
                            if r > 0 and c > 0:
                                row_list.append(cell.text)

                if len(row_list) > 0:
                    if n == 11:
                        bg_list.append(row_list)
                    if n == 12:
                        fb_list.append(row_list)
                    if n == 13:
                        zy_list.append(row_list)
                    if n == 14:
                        ph_list.append(row_list)
                    if n == 15:
                        gz_list.append(row_list)


def write_m_page(n):
    bg = m_prs.slides[n-1]

    for i in bg.shapes:
        if i.has_text_frame and i.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # 本月工作成果
            if n in [14]:
                tf = i.text_frame
                x = 1
                tf.text = ''
                for w in [bg[0] for bg in bg_list]:
                    p = tf.add_paragraph()
                    w = "{}.   {}".format(x, w)
                    p.text = w
                    p.font.size = Pt(14)
                    p.font.name = '宋体-简'
                    p.font.bold = True
                    p.level = 0
                    x += 1
            # 运维工作统计
            if n == 7:
                if i.text_frame.text.find('运维工作统计') < 0:
                    tf = i.text_frame
                    tf.text = "本月主要工作为："
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.size = Pt(12)
                    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    # tf.paragraphs[0].bullet = "●"
                    work_list = []
                    work_list.extend([bg[0] for bg in bg_list])
                    work_list.append("日常配合应用运维或研发人员做问题排查{}次".format(len(ph_list)))
                    work_list.append("支撑发版{}次".format(len(fb_list)))
                    for w in work_list:
                        p = tf.add_paragraph()
                        # p.bullet = "●"
                        p.level = 0
                        p.text = "●  {}".format(w)
                        p.font.color.rgb = RGBColor(255, 0, 0)
                        p.font.size = Pt(10)
                        p.font.name = '微软雅黑'
                        p.font.bold = True
            # es
            if n == 23:
                tf = i.text_frame
                if tf.text.lower().find('es服务个数') >= 0:
                    tf.text = es_des
                    tf.paragraphs[0].font.bold = False
                    tf.paragraphs[0].font.size = Pt(10)
                    tf.paragraphs[0].font.name = '微软雅黑'
                    for x in tf.paragraphs:
                        x.font.bold = False
                        x.font.size = Pt(10)
                        x.font.name = '微软雅黑'

        if i.has_chart:
            c = i.chart
            # 运维工作进展-运维工作统计
            if n == 7:
                # 修改图表chart数据
                chart_data = CategoryChartData()
                chart_data.categories = ['变更', '资源权限管理', '配合操作', '支撑发版', '故障和问题处理']
                chart_data.add_series('系列1', (len(bg_list), len(zy_list), len(ph_list), len(fb_list), len(gz_list)))
                c.replace_data(chart_data)
            # 本月运维情况分析---集群资源可分配容量
            if n == 18:
                title = c.chart_title.text_frame.text
                chart_data = CategoryChartData()
                chart_data.categories = ['外网微服务', '外网后台区', '内网微服务区', '内网后台区']
                if title == 'CPU资源使用情况':
                    chart_data.add_series('已申请', (
                        float('{:.2f}'.format(used(wwwfw_cpu[-1], 'cpu_ysy')/100)),
                        float('{:.2f}'.format(used(wwht_cpu[-1], 'cpu_ysy')/100)),
                        float('{:.2f}'.format(used(nwwfw_cpu[-1], 'cpu_ysy')/100)),
                        float('{:.2f}'.format(used(nwht_cpu[-1], 'cpu_ysy')/100)),
                    ))
                    chart_data.add_series('剩余', (
                        float('{:.2f}'.format(used(wwwfw_cpu[-1], 'cpu_wsy') / 100)),
                        float('{:.2f}'.format(used(wwht_cpu[-1], 'cpu_wsy') / 100)),
                        float('{:.2f}'.format(used(nwwfw_cpu[-1], 'cpu_wsy') / 100)),
                        float('{:.2f}'.format(used(nwht_cpu[-1], 'cpu_wsy') / 100)),
                    ))

                if title == '内存使用情况':
                    chart_data.add_series('已申请', (
                        float('{:.2f}'.format(used(wwwfw_mem[-1], 'mem_ysy')/100)),
                        float('{:.2f}'.format(used(wwht_mem[-1], 'mem_ysy')/100)),
                        float('{:.2f}'.format(used(nwwfw_mem[-1], 'mem_ysy')/100)),
                        float('{:.2f}'.format(used(nwht_mem[-1], 'mem_ysy')/100)),
                    ))
                    chart_data.add_series('剩余', (
                        float('{:.2f}'.format(used(wwwfw_mem[-1], 'mem_wsy') / 100)),
                        float('{:.2f}'.format(used(wwht_mem[-1], 'mem_wsy') / 100)),
                        float('{:.2f}'.format(used(nwwfw_mem[-1], 'mem_wsy') / 100)),
                        float('{:.2f}'.format(used(nwht_mem[-1], 'mem_wsy') / 100)),
                    ))
                c.replace_data(chart_data)
                value_axis = c.value_axis
                tick_labels = value_axis.tick_labels
                tick_labels.number_format = '0%'

            if n == 23:
                chart_data = CategoryChartData()
                chart_data.categories = es_month
                for p in es_index:
                    chart_data.add_series(p[0], p[1])
                c.replace_data(chart_data)

            # 外网微服务区资源使用情况
            dz = {'19': 'wwwfw', '20': 'wwht', '21': 'nwwfw', '22': 'nwht'}
            if n in [19, 20, 21, 22]:
                chart_data = CategoryChartData()
                chart_data.categories = ["{}{}".format(month, i) for i in weeks]
                chart_data.add_series('CPU资源使用率', tuple(
                    [float('{:.2f}'.format(used(i, 'cpu_ysy')/100)) for i in eval('{}_cpu'.format(dz[str(n)]))]
                ))
                chart_data.add_series('内存使用率', tuple(
                    [float('{:.2f}'.format(used(i, 'mem_ysy')/100)) for i in eval('{}_mem'.format(dz[str(n)]))]
                ))
                c.replace_data(chart_data)
                value_axis = c.value_axis
                tick_labels = value_axis.tick_labels
                tick_labels.number_format = '0%'

        if i.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = i.table
            # 写入7.运维工作进展-运维工作统计
            if n == 7:
                tj = [len(bg_list), len(zy_list), len(ph_list), len(fb_list), len(gz_list)]
                for c in range(1, 6):
                    cell = tbl.cell(1, c)
                    cell.text = str(tj[c-1])
                    cell.text_frame.paragraphs[0].font.bold = False
                    cell.text_frame.paragraphs[0].font.name = "微软雅黑"
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

            # 写入8.运维工作进展变更
            if n == 8:
                for r in range(len(bg_list)):
                    cell = tbl.cell(r, 1)
                    cell.text = str(bg_list[r][0])
                    cell.text_frame.paragraphs[0].font.bold = False
                    cell.text_frame.paragraphs[0].font.name = "宋体-简"
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
            # 写入9.运维工作进展-资源权限管理
            if n == 9:
                for r in range(len(zy_list)):
                    for c in range(len(zy_list[0])):
                        cell = tbl.cell(r+1, c+1)
                        cell.text = str(zy_list[r][c])
                        cell.text_frame.paragraphs[0].font.bold = False
                        cell.text_frame.paragraphs[0].font.name = "宋体-简"
                        cell.text_frame.paragraphs[0].font.size = Pt(12)
                        # 居中
                        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            # 写入10.运维工作进展-配合操作及排障
            if n == 10:
                for r in range(len(ph_list)):
                    for c in range(len(ph_list[0])):
                        cell = tbl.cell(r+1, c+1)
                        cell.text = str(ph_list[r][c])
                        cell.text_frame.paragraphs[0].font.bold = False
                        cell.text_frame.paragraphs[0].font.name = "宋体-简"
                        cell.text_frame.paragraphs[0].font.size = Pt(12)
                        # 居中
                        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

            # 写入11.运维工作进展-支撑发版
            if n == 11:
                for r in range(len(fb_list)):
                    for c in range(len(fb_list[0])):
                        cell = tbl.cell(r+1, c)
                        cell.text = str(fb_list[r][c])
                        cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].font.name = "黑体"
                        cell.text_frame.paragraphs[0].font.size = Pt(14)
                        # 居中
                        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            # 写入12.运维工作进展-问题处理
            if n == 12:
                for r in range(len(gz_list)):
                    for c in range(len(gz_list[0])):
                        cell = tbl.cell(r+1, c+1)
                        cell.text = str(gz_list[r][c])
                        cell.text_frame.paragraphs[0].font.bold = False
                        cell.text_frame.paragraphs[0].font.name = "宋体-简"
                        cell.text_frame.paragraphs[0].font.size = Pt(8)
                        # 居中
                        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        m_prs.save(new_m)


def read_latest_zb(n):
    """ 本月最后一周周报 """
    bg = latest_zb_prs.slides[n-1]
    dz = {'19': 'wwwfw', '20': 'wwht', '21': 'nwwfw', '22': 'nwht'}
    for i in bg.shapes:
        if i.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = i.table
            cnt_row, cnt_col = len(tbl.rows), len(tbl.columns)
            for r in range(0, cnt_row):
                row_list = []
                for c in range(0, cnt_col):
                    cell = tbl.cell(r, c)
                    if len(cell.text) != 0:
                        # 19-22慧企运行情况分析–集群资源使用情况
                        if n in [19, 20, 21, 22] and c in [1, 2]:
                            row_list.append(cell.text)
                        if c in [0]:
                            eval("{}_pj".format(dz[str(n)])).append(cell.text)

                eval(dz[str(n)]).append(row_list)

        if i.has_chart:
            c = i.chart
            if n == 23:
                # global es_month, es_index
                es_month.extend(list(c.plots[0].categories))
                es_index.extend([(ii.name, ii.values) for ii in c.series])
                # print("chart", list(c.plots[0].categories))
                # print("chart2", [i.name for i in c.series])
                # print("chart2", [(i.name, i.values) for i in c.series])

        if i.has_text_frame and i.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if n == 24:
                global es_des
                es_des = i.text_frame.text


def used(d, k='cpu_ysy'):
    return d[k]/sum(d.values()) * 100


if __name__ == '__main__':
    main()
    # 慧企运行情况分析–集群资源使用情况
    area = {'wwwfw': '外网微服务', 'wwht': '外网后台', 'nwwfw': '内网微服务', 'nwht': '内网后台'}
    old_wwwfw_pj, old_wwht_pj, old_nwwfw_pj, old_nwht_pj = [], [], [], []
    old_wwwfw, old_wwht, old_nwwfw, old_nwht = [], [], [], []
    page = 19
    prs = Presentation(new_m)
    for fun in ['wwwfw', 'wwht', 'nwwfw', 'nwht']:
        pj = eval('{}_pj'.format(fun))
        rs = eval('{}'.format(fun))
        old_pj = eval('old_{}_pj'.format(fun))
        old_rs = eval('old_{}'.format(fun))

        cpu_used = used(eval("{}_cpu".format(fun))[-1], 'cpu_ysy')
        mem_used = used(eval("{}_mem".format(fun))[-1], 'mem_ysy')
        # print("{}区部署{}个项目，包含{}个服务，运行实例{}个。占用CPU资源升高到{:.1f}%，内存资源升高到{:.1f}%。".format(
        #     area[fun], len(eval('{}_pj'.format(fun))), len(eval(fun)), sum([int(i[1]) for i in eval(fun)]),
        #     cpu_used, mem_used
        # ))
        cpu = eval("{}_cpu".format(fun))[-1]
        mem = eval("{}_mem".format(fun))[-1]
        des = "{}区CPU使用率{:.0f}%，共计使用{}/{}核。内存{:.0f}%，使用{}/{}GB。".format(
            area[fun], cpu_used, int(cpu['cpu_ysy']), int(sum(cpu.values())),
            mem_used, int(mem['mem_ysy']), int(sum(mem.values()))
        )
        bg = prs.slides[page - 1]
        for i in bg.shapes:
            # print(page, i.shape_type)
            if i.has_text_frame and i.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                tf = i.text_frame
                tf.text = des
                tf.paragraphs[0].font.bold = False
                tf.paragraphs[0].font.size = Pt(10)

            if i.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = i.table
                cnt_row, cnt_col = len(tbl.rows), len(tbl.columns)
                for r in range(cnt_row):
                    row_list = []
                    for c in range(0, cnt_col):
                        cell = tbl.cell(r, c)
                        if len(cell.text) != 0:
                            # 19-22慧企运行情况分析–集群资源使用情况
                            if page in [19, 20, 21, 22] and c in [1, 2]:
                                row_list.append(cell.text)
                            if page in [19, 20, 21, 22] and c in [0]:
                                old_pj.append(cell.text)

                    old_rs.append(row_list)

        page += 1
        # 服务项目变化对比
        print(area[fun].center(50, '-'))
        pj = [i.strip().lower() for i in pj]
        old_pj = [i.strip().lower() for i in old_pj]
        print("新增项目", set(pj) - set(old_pj))
        print("减去项目", set(old_pj) - set(pj))
        app = [i[0].strip().lower() for i in rs if len(i) == 2 and i[1].isdigit() and not i[0].isdigit()]
        old_app = [i[0].strip().lower() for i in old_rs if len(i) == 2 and i[1].isdigit() and not i[0].isdigit()]
        add_app = set(app) - set(old_app)
        for a in add_app:
            num = [i[1] for i in rs if len(i) == 2 and i[0].strip().lower() == a][0]
            print("新增应用{}, 实例数{}".format(a, num))
        print("减去应用", set(old_app) - set(app))
        for x in [x for x in rs if len(x) == 2 and x[1].isdigit() and not x[0].isdigit()]:
            for y in [y for y in old_rs if len(y) == 2 and y[1].isdigit() and not y[0].isdigit()]:
                if y[0].strip() == x[0].strip() and y[1] != x[1]:
                    print("{}应用实例发生变化：由{}变为{}.".format(y[0], y[1], x[1]))
    prs.save(new_m)





